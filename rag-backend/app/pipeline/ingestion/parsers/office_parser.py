from __future__ import annotations

import asyncio
from pathlib import Path

from app.observability.logger import get_logger
from app.pipeline.ingestion.parsers.base import BaseParser, ParsedDocument, ParsedElement

logger = get_logger(__name__)


class OfficeParser(BaseParser):
    """Word / Excel / PPT 解析器，使用 Unstructured"""

    parser_name = "office_parser"

    async def parse(self, file_path: Path) -> ParsedDocument:
        suffix = file_path.suffix.lower()
        if suffix == ".docx":
            return await self._parse_docx(file_path)
        elif suffix in (".xlsx", ".xls"):
            return await self._parse_excel(file_path)
        elif suffix in (".pptx", ".ppt"):
            return await self._parse_pptx(file_path)
        else:
            raise ValueError(f"OfficeParser does not support {suffix}")

    async def _parse_docx(self, file_path: Path) -> ParsedDocument:
        from unstructured.partition.docx import partition_docx

        def _sync():
            return partition_docx(filename=str(file_path), include_page_breaks=True)

        raw_elements = await asyncio.to_thread(_sync)
        return self._convert_elements(raw_elements, "unstructured_docx")

    async def _parse_excel(self, file_path: Path) -> ParsedDocument:
        """Excel：每个 Sheet 的表格作为独立 table chunk"""
        import openpyxl

        def _sync():
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            elements = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                # 转换为 Markdown 表格
                md_rows = []
                for i, row in enumerate(rows):
                    cells = [str(c) if c is not None else "" for c in row]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                content = f"**Sheet: {sheet_name}**\n\n" + "\n".join(md_rows)
                elements.append(
                    ParsedElement(
                        content=content,
                        element_type="table",
                        metadata={"sheet_name": sheet_name},
                        confidence=0.95,
                    )
                )
            return elements

        elements = await asyncio.to_thread(_sync)
        return ParsedDocument(elements=elements, parser_used="openpyxl")

    async def _parse_pptx(self, file_path: Path) -> ParsedDocument:
        """PPT：每页幻灯片作为独立 chunk"""
        from pptx import Presentation

        def _sync():
            prs = Presentation(str(file_path))
            elements = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                title = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                if shape.shape_type == 13 or (hasattr(shape, "name") and "title" in shape.name.lower()):
                                    title = text
                                else:
                                    texts.append(text)

                content = (f"# {title}\n\n" if title else f"# Slide {slide_num}\n\n") + "\n".join(texts)
                if content.strip():
                    elements.append(
                        ParsedElement(
                            content=content,
                            element_type="text",
                            page_number=slide_num,
                            section_path=[title] if title else [f"Slide {slide_num}"],
                            confidence=0.9,
                        )
                    )
            return elements, len(prs.slides)

        elements, page_count = await asyncio.to_thread(_sync)
        return ParsedDocument(elements=elements, page_count=page_count, parser_used="python_pptx")

    def _convert_elements(self, raw_elements: list, parser_used: str) -> ParsedDocument:
        parsed_elements = []
        current_section: list[str] = []

        for el in raw_elements:
            el_type = type(el).__name__.lower()
            page_num = getattr(el.metadata, "page_number", None)

            if "table" in el_type:
                chunk_type = "table"
            elif "title" in el_type:
                chunk_type = "heading"
                if el.text.strip():
                    current_section = [el.text.strip()]
            else:
                chunk_type = "text"

            content = el.text if hasattr(el, "text") else str(el)
            if not content.strip():
                continue

            parsed_elements.append(
                ParsedElement(
                    content=content,
                    element_type=chunk_type,
                    page_number=page_num,
                    section_path=list(current_section),
                    confidence=0.9,
                )
            )

        return ParsedDocument(elements=parsed_elements, parser_used=parser_used)
