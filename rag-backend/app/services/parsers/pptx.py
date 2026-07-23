from io import BytesIO
from pathlib import Path

from pptx import Presentation


def parse_pptx(source: bytes | str | Path) -> list[dict]:
    """Returns list of {page_num=<slide#>, text, kind='slide'}.

    每个 slide 一个 chunk；包括 shape 文本 + 演讲者备注。
    """
    if isinstance(source, (bytes, bytearray)):
        prs = Presentation(BytesIO(source))
    else:
        prs = Presentation(str(source))

    chunks: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text or "").strip()
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                parts.append(f"[备注] {notes}")
        text = "\n".join(parts).strip()
        if text:
            chunks.append({
                "page_num": i, "text": text, "kind": "slide",
                "parse_confidence": 0.95,
                "section_path": [title] if title else [],
            })
    return chunks
